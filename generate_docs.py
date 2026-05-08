"""
generate_docs.py — Create PowerPoint presentation and Word document from simulation report.
Includes methodology, sensitivity analysis, break-even, and references.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

REPORT_DIR = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "output", "report_20260428_113016"
)

PLOTS = {
    "passenger_metrics": os.path.join(REPORT_DIR, "passenger_metrics.png"),
    "trip_comparison": os.path.join(REPORT_DIR, "trip_comparison.png"),
    "fundamental_diagrams": os.path.join(REPORT_DIR, "fundamental_diagrams.png"),
    "space_time": os.path.join(REPORT_DIR, "space_time_diagram.png"),
    "sensitivity_demand": os.path.join(REPORT_DIR, "sensitivity_demand.png"),
    "sensitivity_congestion": os.path.join(REPORT_DIR, "sensitivity_congestion.png"),
    "breakeven_heatmap": os.path.join(REPORT_DIR, "breakeven_heatmap.png"),
    "mode_choice": os.path.join(REPORT_DIR, "mode_choice.png"),
}

# ── Color palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT     = RGBColor(0x00, 0xCC, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
TABLE_HDR  = RGBColor(0x00, 0x88, 0xBB)


def _set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def _add_bullet(tf, text, font_size=16, color=WHITE, bold=False, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    return p


def _add_img(slide, key, left, top, width, height):
    path = PLOTS.get(key, "")
    if path and os.path.isfile(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))


def _make_table(slide, rows_data, left, top, width, height):
    """Create a styled table on a slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    tbl = tbl_shape.table
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if ri == 0:
                p.font.bold = True
                cell.fill.fore_color.rgb = TABLE_HDR
            else:
                cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x40)
    return tbl


# ═══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════════

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 1. Title ─────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 1.0, 1.5, 11.3, 1.5, "SUMO-AAMSim",
                  font_size=48, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, 1.0, 3.0, 11.3, 1.0,
                  "eVTOL Ground-Traffic Integration Simulation",
                  font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, 1.0, 4.5, 11.3, 0.6,
                  "Tampa Bay AAM Study  |  Tampa <-> Brandon Corridor",
                  font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, 1.0, 5.5, 11.3, 0.5,
                  "Marvin Osei-Kuffour  |  University of South Florida  |  April 2026",
                  font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── 2. Methodology ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.3, 12.3, 0.8,
                  "Methodology", font_size=36, bold=True, color=ACCENT)

    tf = _add_text_box(slide, 0.5, 1.2, 6.0, 5.8,
                       "Simulation Architecture", font_size=20, bold=True, color=ACCENT)
    for item in [
        "Multi-modal co-simulation framework:",
        "  SUMO: microscopic ground traffic (TraCI API)",
        "  CarlaAir: 3D visualization (CARLA + AirSim)",
        "  Python orchestrator: real-time bridge",
        "",
        "Ground traffic: Krauss car-following model",
        "  2 lanes/dir, 50 km/h, v/c = 0.92 peak hour",
        "  4,405 veh/hr/dir background flow",
        "",
        "eVTOL flights: AirSim multirotor drone",
        "  Pose interpolation at 100 km/h cruise",
        "  Per-vertiport terrain-aware landing",
    ]:
        _add_bullet(tf, item, font_size=13, color=WHITE)

    tf2 = _add_text_box(slide, 6.5, 1.2, 6.3, 5.8,
                        "Data-Driven Passenger Model", font_size=20, bold=True, color=ACCENT)
    for item in [
        "Tampa Bay AAM study Excel data:",
        "  168 pax Tampa->Brandon (VP-Alpha)",
        "  148 pax Brandon->Tampa (VP-Beta)",
        "  Actual arrival times over 5-hour peak period",
        "",
        "Boarding rules with timeout triggers:",
        "  30 min renege, 15 min solo, 10 min 2nd pax",
        "  4 pax capacity -> immediate takeoff",
        "",
        "CarlaAir eVTOL integration:",
        "  CARLA (port 2000): ground vehicle mirroring",
        "  AirSim (port 41451): drone flight control",
        "  SUMO eVTOL route -> AirSim pose interpolation",
        "  Pedestrians spawn at vertiports on landing",
    ]:
        _add_bullet(tf2, item, font_size=13, color=WHITE)

    # ── 3. Simulation Overview ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.3, 12.3, 0.8,
                  "Simulation Parameters", font_size=36, bold=True, color=ACCENT)

    tf = _add_text_box(slide, 0.8, 1.3, 5.5, 5.5,
                       "Parameters", font_size=22, bold=True, color=ACCENT)
    for item in [
        "Duration: 5 hours (18,000 steps)",
        "Road: 2 lanes/direction, 50 km/h",
        "Volume-to-capacity ratio: 0.92 (peak hour)",
        "Background traffic: 4,405 veh/hr/direction",
        "eVTOL cruise speed: 100 km/h",
        "eVTOL capacity: 4 passengers",
        "VP distance: 30 km (Tampa <-> Brandon)",
        "Passenger data: Tampa Bay AAM study",
    ]:
        _add_bullet(tf, item, font_size=16, color=WHITE)

    tf2 = _add_text_box(slide, 6.8, 1.3, 5.8, 5.5,
                        "Boarding Rules", font_size=22, bold=True, color=ACCENT)
    for item in [
        "1. Passenger reneges after 30 min wait",
        "2. First pax alone 15 min -> takeoff",
        "3. 2+ pax, 2nd waited 10 min -> takeoff",
        "4. 4 passengers (full) -> immediate takeoff",
    ]:
        _add_bullet(tf2, item, font_size=16, color=WHITE)
    _add_bullet(tf2, "", font_size=10)
    _add_bullet(tf2, "Two vertiports:", font_size=18, bold=True, color=ACCENT)
    for item in [
        "VP-Alpha (Tampa): 168 scheduled passengers",
        "VP-Beta (Brandon): 148 scheduled passengers",
        "Total: 316 passengers over 5 hours",
    ]:
        _add_bullet(tf2, item, font_size=16, color=WHITE)

    # ── 4. Key Results ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.3, 12.3, 0.8,
                  "Key Results", font_size=36, bold=True, color=ACCENT)
    for value, label, left in [
        ("314", "Passengers\nServed", 0.5),
        ("80", "Total\nFlights", 3.5),
        ("0.0%", "Reneging\nRate", 6.5),
        ("~31 min", "Avg Time\nSaved", 9.5),
    ]:
        _add_text_box(slide, left, 1.5, 2.8, 1.2, value,
                      font_size=44, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, left, 2.8, 2.8, 1.0, label,
                      font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    _make_table(slide, [
        ["Vertiport", "Passengers", "Served", "Reneged",
         "Avg OVWT", "Avg IVWT", "eVTOL Trip", "Time Saved"],
        ["VP-Alpha (Tampa)", "168", "166", "0",
         "0.0 min", "3.0 min", "41.0 min", "31.0 min"],
        ["VP-Beta (Brandon)", "148", "148", "0",
         "0.0 min", "2.8 min", "40.8 min", "31.2 min"],
    ], 0.5, 4.2, 12.3, 2.2)

    # ── 5. Passenger Metrics Plot ────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Passenger Metrics", font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "passenger_metrics", 0.8, 1.0, 11.7, 6.0)

    # ── 6. Trip Comparison ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Trip Time Comparison: Ground vs eVTOL",
                  font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "trip_comparison", 1.5, 1.0, 10.3, 6.0)
    _add_text_box(slide, 0.5, 6.8, 12.3, 0.5,
                  "Ground: 72 min (30 km @ 25 km/h)  |  "
                  "eVTOL: ~41 min (taxi+wait+flight+taxi)  |  "
                  "Savings: ~31 min/pax",
                  font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── 7. Fundamental Diagrams ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Fundamental Diagrams", font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "fundamental_diagrams", 0.8, 1.0, 11.7, 6.0)

    # ── 8. Space-Time Diagram ────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Space-Time Traffic Diagram", font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "space_time", 0.8, 1.0, 11.7, 6.0)

    # ── 9. Sensitivity: Demand ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Sensitivity Analysis: Demand Variation",
                  font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "sensitivity_demand", 0.5, 1.0, 12.3, 6.0)

    # ── 10. Sensitivity: Congestion ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Sensitivity Analysis: Traffic Congestion",
                  font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "sensitivity_congestion", 0.3, 1.0, 12.7, 5.5)

    # ── 11. Break-Even Distance ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Break-Even Distance: eVTOL vs Ground",
                  font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "breakeven_heatmap", 0.8, 1.0, 11.7, 6.0)

    # ── 12. Mode Choice ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.2, 12.3, 0.7,
                  "Mode Choice: When to Choose eVTOL?",
                  font_size=32, bold=True, color=ACCENT)
    _add_img(slide, "mode_choice", 0.5, 1.0, 12.3, 6.0)

    # ── 13. Key Findings from Sensitivity ────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.3, 12.3, 0.8,
                  "Sensitivity Analysis: Key Findings",
                  font_size=36, bold=True, color=ACCENT)

    _make_table(slide, [
        ["v/c Ratio", "Ground Speed", "Ground Trip", "eVTOL Trip",
         "Time Saved", "Break-Even Dist."],
        ["0.50 (free)", "38.0 km/h", "47 min", "41 min", "6 min", "23.5 km"],
        ["0.70 (moderate)", "33.2 km/h", "54 min", "41 min", "13 min", "19.1 km"],
        ["0.85 (heavy)", "27.7 km/h", "65 min", "41 min", "24 min", "14.7 km"],
        ["0.92 (peak)", "24.9 km/h", "72 min", "41 min", "31 min", "12.7 km"],
        ["1.00 (saturated)", "21.6 km/h", "83 min", "41 min", "42 min", "10.6 km"],
    ], 0.8, 1.5, 11.7, 3.5)

    tf = _add_text_box(slide, 0.8, 5.3, 11.5, 2.0, "", font_size=15, color=WHITE)
    _add_bullet(tf, "eVTOL advantage increases sharply with congestion (v/c > 0.85)",
                font_size=15, color=WHITE)
    _add_bullet(tf, "Break-even distance drops from 23.5 km (free flow) to 10.6 km (saturated)",
                font_size=15, color=WHITE)
    _add_bullet(tf, "For Tampa-Brandon (30 km), eVTOL is faster at all congestion levels above v/c = 0.50",
                font_size=15, color=WHITE)
    _add_bullet(tf, "Higher demand -> shorter IVWT -> more time savings per passenger",
                font_size=15, color=WHITE)

    # ── 14. Conclusions ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.3, 12.3, 0.8,
                  "Conclusions", font_size=36, bold=True, color=ACCENT)

    tf = _add_text_box(slide, 0.8, 1.3, 11.5, 5.5, "", font_size=18, color=WHITE)
    for title, desc in [
        ("Time Savings",
         "eVTOL saves ~31 min/trip at peak hour (72 min ground vs ~41 min eVTOL). "
         "Savings increase to 42 min under saturated conditions (v/c = 1.0)."),
        ("Break-Even Distance",
         "eVTOL becomes advantageous for trips > 12.7 km at peak hour (v/c = 0.92). "
         "In free-flow conditions, trips must exceed 23.5 km. Tampa-Brandon (30 km) "
         "benefits at all congestion levels above v/c = 0.50."),
        ("System Efficiency",
         "0% reneging, avg 3.9 pax/flight, zero spillback. Demand-supply is well matched."),
        ("Mode Choice Guidance",
         "Choose eVTOL when: trip distance > break-even AND congestion is moderate-to-heavy. "
         "Ground transport is preferable for short trips (< 15 km) or free-flow conditions."),
    ]:
        _add_bullet(tf, title, font_size=18, bold=True, color=ACCENT)
        _add_bullet(tf, desc, font_size=14, color=WHITE)
        _add_bullet(tf, "", font_size=6)

    # ── 15. References ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 0.5, 0.3, 12.3, 0.8,
                  "References", font_size=36, bold=True, color=ACCENT)

    tf = _add_text_box(slide, 0.8, 1.3, 11.5, 5.5, "", font_size=14, color=WHITE)
    refs = [
        "[1] Dosovitskiy, A. et al. (2017). CARLA: An Open Urban Driving Simulator. "
        "Conference on Robot Learning (CoRL), pp. 1-16.",
        "[2] Shah, S. et al. (2018). AirSim: High-Fidelity Visual and Physical Simulation "
        "for Autonomous Vehicles. Field and Service Robotics, pp. 621-635.",
        "[3] Lopez, P.A. et al. (2018). Microscopic Traffic Simulation using SUMO. "
        "IEEE Intelligent Transportation Systems Conference (ITSC), pp. 2575-2582.",
        "[4] Balasubramaniam, S. et al. (2024). CarlaAir: A Unified Platform for "
        "Air-Ground Co-Simulation of Urban Air Mobility. Preprint / Working Paper.",
        "[5] Tampa Bay Regional Planning Council (2024). Tampa Bay Advanced Air Mobility "
        "(AAM) Feasibility Study. Passenger demand data for Tampa-Brandon corridor.",
        "[6] Bureau of Public Roads (1964). Traffic Assignment Manual. "
        "U.S. Dept. of Commerce. (BPR speed-flow function.)",
        "[7] Krauss, S. (1998). Microscopic Modeling of Traffic Flow: Investigation of "
        "Collision Free Vehicle Dynamics. PhD Thesis, DLR.",
        "[8] FAA (2023). Advanced Air Mobility (AAM) Implementation Plan. "
        "Federal Aviation Administration.",
    ]
    for ref in refs:
        _add_bullet(tf, ref, font_size=12, color=LIGHT_GRAY)

    # ── 16. Thank You ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_text_box(slide, 1.0, 2.5, 11.3, 1.5, "Thank You",
                  font_size=48, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, 1.0, 4.2, 11.3, 1.0,
                  "SUMO-AAMSim  |  eVTOL Ground-Traffic Integration",
                  font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, 1.0, 5.2, 11.3, 0.6,
                  "Marvin Osei-Kuffour  |  University of South Florida",
                  font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    out_path = os.path.join(REPORT_DIR, "SUMO_AAMSim_Report.pptx")
    prs.save(out_path)
    print(f"PowerPoint saved: {out_path}")
    return out_path


# ═══════════════════════════════════════════════════════════════════════════════
#  WORD DOCUMENT
# ═══════════════════════════════════════════════════════════════════════════════

def create_docx():
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)

    # ── Title page ───────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("SUMO-AAMSim Simulation Report")
    run.bold = True
    run.font.size = DocxPt(26)
    run.font.color.rgb = DocxRGB(0x00, 0x88, 0xBB)

    for text, size, color in [
        ("eVTOL Ground-Traffic Integration Simulation", 14, 0x555555),
        ("Marvin Osei-Kuffour  |  University of South Florida  |  April 2026", 11, 0x777777),
        ("Generated: 2026-04-28", 10, 0x999999),
    ]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.size = DocxPt(size)
        run.font.color.rgb = DocxRGB(color >> 16, (color >> 8) & 0xFF, color & 0xFF)

    # ── 1. Methodology ───────────────────────────────────────────────────────
    doc.add_heading("1. Methodology", level=1)

    doc.add_heading("1.1 Simulation Architecture", level=2)
    doc.add_paragraph(
        "This study employs a multi-modal co-simulation framework integrating three "
        "platforms to evaluate eVTOL ground-traffic interactions during peak-hour conditions "
        "on the Tampa-Brandon corridor:"
    )
    methods = [
        ("SUMO (Simulation of Urban Mobility)",
         "Provides microscopic ground traffic simulation using the Krauss car-following model. "
         "Controlled via the TraCI (Traffic Control Interface) API for real-time vehicle "
         "injection, routing, and metric collection. Background traffic flows at 4,405 veh/hr/dir "
         "to achieve a volume-to-capacity ratio of 0.92."),
        ("CarlaAir (CARLA + AirSim)",
         "A unified Unreal Engine 4.26 platform that runs CARLA (port 2000) for ground vehicle "
         "rendering and AirSim (port 41451) for multirotor drone flight simulation in a single "
         "process. This enables synchronized air-ground visualization without separate processes."),
        ("Python Orchestrator",
         "A real-time bridge that mirrors SUMO ground vehicles into CARLA 3D space, translates "
         "eVTOL departures from the passenger queue system into AirSim drone flights using pose "
         "interpolation, and spawns CARLA pedestrians at vertiports when drones land."),
    ]
    for title, desc in methods:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(f"{title}: ")
        run.bold = True
        p.add_run(desc)

    doc.add_heading("1.2 eVTOL Integration in CarlaAir", level=2)
    doc.add_paragraph(
        "The eVTOL aircraft is simulated as an AirSim multirotor drone within the CarlaAir "
        "environment. Key implementation details:"
    )
    evtol_details = [
        "The drone controller uses simSetVehiclePose() for direct position interpolation "
        "rather than async movement commands, ensuring reliable flight paths.",
        "Flight phases: IDLE -> FLYING (tracking SUMO eVTOL route) -> DESCEND -> LANDED "
        "(5s ground pause for passenger visualization) -> RETURN_FLY -> RETURN_LAND -> IDLE.",
        "Per-vertiport terrain heights are measured via CARLA ray-casting to ensure the drone "
        "lands at the correct ground elevation at each vertiport (VP-Alpha: 0.0m, VP-Beta: 38.64m).",
        "Cruise altitude is set 25m above the higher vertiport to maintain clearance over "
        "varying terrain between Tampa and Brandon.",
        "Coordinate transforms: SUMO -> CARLA (y-flip from OpenDRIVE) -> AirSim NED "
        "(measured offsets: dx=172.2, dy=-183.9, dz=27.5).",
    ]
    for detail in evtol_details:
        doc.add_paragraph(detail, style="List Bullet")

    doc.add_heading("1.3 Data-Driven Passenger Model", level=2)
    doc.add_paragraph(
        "Passenger arrival times are sourced from the Tampa Bay AAM Feasibility Study Excel "
        "dataset, which provides individual passenger arrival timestamps for the Tampa-Brandon "
        "corridor. The dataset contains 168 passengers departing Tampa (VP-Alpha) and 148 "
        "passengers departing Brandon (VP-Beta) over a 5-hour peak period."
    )
    doc.add_paragraph(
        "A queuing model with the following boarding rules governs eVTOL dispatch:"
    )
    rules = [
        "Passenger reneges (leaves queue) after 30 minutes of waiting at the vertiport.",
        "If the first boarded passenger waits 15 minutes with no second passenger arriving, "
        "the eVTOL takes off with a single passenger.",
        "If 2+ passengers are boarded and the second passenger has waited more than 10 minutes, "
        "the eVTOL takes off with the current load.",
        "When the eVTOL reaches full capacity (4 passengers), it takes off immediately.",
    ]
    for rule in rules:
        doc.add_paragraph(rule, style="List Number")

    doc.add_heading("1.4 Speed-Flow Model", level=2)
    doc.add_paragraph(
        "Ground vehicle travel times are modeled using the Bureau of Public Roads (BPR) "
        "speed-flow function, calibrated to match the SUMO simulation results:"
    )
    doc.add_paragraph(
        "    v = v_free / (1 + 0.85 * (v/c)^4)"
    )
    doc.add_paragraph(
        "Where v_free = 40 km/h (urban free-flow including signals), and v/c is the "
        "volume-to-capacity ratio. At the simulated peak hour (v/c = 0.92), this yields "
        "an effective speed of ~25 km/h, matching the 72-minute ground trip time for 30 km "
        "observed in the SUMO simulation."
    )

    # ── 2. Simulation Parameters ─────────────────────────────────────────────
    doc.add_heading("2. Simulation Parameters", level=1)
    params = [
        ("Simulation Duration", "5 hours (18,000 simulation steps)"),
        ("Road Configuration", "2 lanes per direction, 50 km/h speed limit"),
        ("Volume-to-Capacity Ratio", "0.92 (peak hour)"),
        ("Background Traffic", "4,405 vehicles/hr/direction"),
        ("eVTOL Cruise Speed", "100 km/h"),
        ("eVTOL Capacity", "4 passengers"),
        ("Vertiport Distance", "30 km (Tampa <-> Brandon)"),
        ("First/Last Mile Taxi", "5 km at 30 km/h average"),
        ("Passenger Data Source", "Tampa Bay AAM Feasibility Study"),
    ]
    for label, value in params:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(f"{label}: ")
        run.bold = True
        p.add_run(value)

    # ── 3. Key Results ───────────────────────────────────────────────────────
    doc.add_heading("3. Key Results", level=1)

    doc.add_heading("3.1 Overall Summary", level=2)
    tbl = doc.add_table(rows=5, cols=2, style="Light Shading Accent 1")
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    for ri, (label, val) in enumerate([
        ("Total Passengers", "316"),
        ("Passengers Served", "314"),
        ("Passengers Reneged", "0"),
        ("Reneging Rate", "0.0%"),
        ("Total Flights", "80"),
    ]):
        tbl.cell(ri, 0).text = label
        tbl.cell(ri, 1).text = val
        for ci in range(2):
            for p in tbl.cell(ri, ci).paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("")

    doc.add_heading("3.2 Per-Vertiport Breakdown", level=2)
    tbl2 = doc.add_table(rows=3, cols=8, style="Light Shading Accent 1")
    tbl2.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ["Vertiport", "Passengers", "Served", "Reneged",
               "Avg OVWT", "Avg IVWT", "eVTOL Trip", "Time Saved"]
    rows = [
        ["VP-Alpha\n(Tampa)", "168", "166", "0", "0.0 min", "3.0 min", "41.0 min", "31.0 min"],
        ["VP-Beta\n(Brandon)", "148", "148", "0", "0.0 min", "2.8 min", "40.8 min", "31.2 min"],
    ]
    for ci, hdr in enumerate(headers):
        cell = tbl2.cell(0, ci)
        cell.text = hdr
        for p in cell.paragraphs:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in p.runs:
                run.bold = True
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            tbl2.cell(ri, ci).text = val
            for p in tbl2.cell(ri, ci).paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("3.3 Key Definitions", level=2)
    for term, defn in [
        ("Out-of-Vehicle Wait Time (OVWT)", "Time from arrival at vertiport to boarding."),
        ("In-Vehicle Wait Time (IVWT)", "Time from boarding to takeoff."),
        ("Reneging Rate", "Passengers who left / total arrivals."),
        ("eVTOL Trip Time", "Taxi to VP + OVWT + IVWT + Flight + Taxi from VP."),
        ("Time Saved", "Ground trip time minus eVTOL total trip time."),
    ]:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(f"{term}: ")
        run.bold = True
        p.add_run(defn)

    # ── 4. Plots ─────────────────────────────────────────────────────────────
    doc.add_heading("4. Passenger Metrics", level=1)
    _add_docx_img(doc, "passenger_metrics", 6.5)

    doc.add_heading("5. Trip Time Comparison", level=1)
    doc.add_paragraph(
        "Ground transport: 72 min (30 km at 25 km/h peak hour). "
        "eVTOL total trip: ~41 min (10 min taxi + 0 min OVWT + 3 min IVWT + 18 min flight + 10 min taxi). "
        "Time saved: ~31 min per passenger."
    )
    _add_docx_img(doc, "trip_comparison", 6.0)

    doc.add_heading("6. Fundamental Diagrams", level=1)
    doc.add_paragraph(
        "Flow-density, speed-density, and speed-flow relationships on the main corridor "
        "during 5-hour peak simulation (v/c = 0.92)."
    )
    _add_docx_img(doc, "fundamental_diagrams", 6.5)

    doc.add_heading("7. Space-Time Traffic Diagram", level=1)
    doc.add_paragraph(
        "Space-time diagrams show traffic wave propagation on eastbound/westbound corridors. "
        "Color indicates speed. Congestion clusters appear near vertiport approach junctions."
    )
    _add_docx_img(doc, "space_time", 6.5)

    # ── 8. Sensitivity Analysis ──────────────────────────────────────────────
    doc.add_heading("8. Sensitivity Analysis", level=1)

    doc.add_heading("8.1 Demand Variation", level=2)
    doc.add_paragraph(
        "Passenger demand was varied from 0.25x to 3.0x the base rate (63.2 pax/hr) while "
        "holding traffic congestion constant at v/c = 0.92. Higher demand reduces in-vehicle "
        "wait time (flights fill faster) and increases time savings per passenger."
    )
    _add_docx_img(doc, "sensitivity_demand", 6.5)

    doc.add_heading("8.2 Traffic Congestion Variation", level=2)
    doc.add_paragraph(
        "The volume-to-capacity ratio was varied from 0.30 to 1.00, representing conditions "
        "from free-flow to fully saturated. eVTOL time savings increase dramatically with "
        "congestion because eVTOL cruise speed (100 km/h) is unaffected by ground traffic."
    )
    _add_docx_img(doc, "sensitivity_congestion", 6.5)

    doc.add_heading("8.3 Break-Even Distance", level=2)
    doc.add_paragraph(
        "The break-even distance is the minimum trip distance at which eVTOL becomes faster "
        "than ground transport. It accounts for the fixed overhead of eVTOL travel (taxi to/from "
        "vertiport, waiting time) versus the distance-dependent speed advantage of eVTOL."
    )
    # Break-even table
    tbl3 = doc.add_table(rows=6, cols=6, style="Light Shading Accent 1")
    tbl3.alignment = WD_TABLE_ALIGNMENT.CENTER
    be_data = [
        ["v/c Ratio", "Ground Speed", "Ground Trip", "eVTOL Trip", "Time Saved", "Break-Even"],
        ["0.50", "38.0 km/h", "47 min", "41 min", "6 min", "23.5 km"],
        ["0.70", "33.2 km/h", "54 min", "41 min", "13 min", "19.1 km"],
        ["0.85", "27.7 km/h", "65 min", "41 min", "24 min", "14.7 km"],
        ["0.92", "24.9 km/h", "72 min", "41 min", "31 min", "12.7 km"],
        ["1.00", "21.6 km/h", "83 min", "41 min", "42 min", "10.6 km"],
    ]
    for ri, row in enumerate(be_data):
        for ci, val in enumerate(row):
            tbl3.cell(ri, ci).text = val
            for p in tbl3.cell(ri, ci).paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if ri == 0:
                    for run in p.runs:
                        run.bold = True

    doc.add_paragraph("")
    _add_docx_img(doc, "breakeven_heatmap", 6.0)

    doc.add_heading("8.4 Mode Choice Analysis", level=2)
    doc.add_paragraph(
        "The mode choice analysis compares door-to-door trip times for ground transport and "
        "eVTOL across different distances and congestion levels. Key findings:"
    )
    findings = [
        "At peak hour (v/c = 0.92), eVTOL is faster for trips exceeding 12.7 km.",
        "In free-flow conditions (v/c = 0.50), trips must exceed 23.5 km for eVTOL advantage.",
        "For the Tampa-Brandon corridor (30 km), eVTOL is advantageous at all congestion "
        "levels above v/c = 0.50.",
        "Short urban trips (< 15 km) generally favor ground transport due to the fixed "
        "overhead of eVTOL (taxi to/from vertiport + waiting).",
        "The first/last mile taxi component (5 km each at 30 km/h = 10 min each) is the "
        "dominant overhead that sets the break-even floor.",
    ]
    for f in findings:
        doc.add_paragraph(f, style="List Bullet")
    _add_docx_img(doc, "mode_choice", 6.5)

    # ── 9. Conclusions ───────────────────────────────────────────────────────
    doc.add_heading("9. Conclusions", level=1)
    conclusions = [
        ("Time Savings",
         "eVTOL passengers save approximately 31 minutes per trip at peak hour (v/c = 0.92), "
         "increasing to 42 minutes under saturated conditions (v/c = 1.0). The savings arise "
         "primarily from the eVTOL's congestion-independent cruise speed of 100 km/h."),
        ("Break-Even Distance",
         "The minimum trip distance for eVTOL advantage ranges from 10.6 km (saturated) to "
         "23.5 km (free-flow). For the Tampa-Brandon corridor (30 km), eVTOL is faster at "
         "all congestion levels above v/c = 0.50."),
        ("System Efficiency",
         "0% passenger reneging rate and near-full flights (avg 3.9 pax/flight) demonstrate "
         "effective demand-supply matching. Zero spillback episodes confirm adequate vertiport "
         "parking capacity."),
        ("Mode Choice Guidance",
         "eVTOL is recommended for trips exceeding the break-even distance during moderate-to-heavy "
         "congestion. Ground transport remains preferable for short trips (< 15 km) or in "
         "free-flow conditions where the eVTOL taxi overhead exceeds the speed advantage."),
        ("Scalability",
         "The data-driven passenger queue model provides a foundation for evaluating additional "
         "vertiport pairs, demand scenarios, and fleet sizing strategies."),
    ]
    for title, desc in conclusions:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(f"{title}: ")
        run.bold = True
        p.add_run(desc)

    # ── 10. References ───────────────────────────────────────────────────────
    doc.add_heading("10. References", level=1)
    refs = [
        "Dosovitskiy, A., Ros, G., Codevilla, F., Lopez, A., & Koltun, V. (2017). "
        "CARLA: An Open Urban Driving Simulator. Conference on Robot Learning (CoRL), pp. 1-16.",
        "Shah, S., Dey, D., Lovett, C., & Kapoor, A. (2018). AirSim: High-Fidelity Visual "
        "and Physical Simulation for Autonomous Vehicles. Field and Service Robotics, pp. 621-635.",
        "Lopez, P.A., Behrisch, M., Bieker-Walz, L., Erdmann, J., et al. (2018). Microscopic "
        "Traffic Simulation using SUMO. IEEE Intelligent Transportation Systems Conference (ITSC).",
        "Balasubramaniam, S. et al. (2024). CarlaAir: A Unified Platform for Air-Ground "
        "Co-Simulation of Urban Air Mobility. Working Paper.",
        "Tampa Bay Regional Planning Council (2024). Tampa Bay Advanced Air Mobility (AAM) "
        "Feasibility Study. Passenger demand data for Tampa-Brandon corridor.",
        "Bureau of Public Roads (1964). Traffic Assignment Manual. U.S. Dept. of Commerce.",
        "Krauss, S. (1998). Microscopic Modeling of Traffic Flow: Investigation of Collision "
        "Free Vehicle Dynamics. PhD Thesis, DLR.",
        "Federal Aviation Administration (2023). Advanced Air Mobility (AAM) Implementation Plan.",
    ]
    for i, ref in enumerate(refs, 1):
        doc.add_paragraph(f"[{i}] {ref}")

    out_path = os.path.join(REPORT_DIR, "SUMO_AAMSim_Report.docx")
    doc.save(out_path)
    print(f"Word document saved: {out_path}")
    return out_path


def _add_docx_img(doc, key, width_in):
    path = PLOTS.get(key, "")
    if path and os.path.isfile(path):
        doc.add_picture(path, width=DocxInches(width_in))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER


if __name__ == "__main__":
    create_pptx()
    create_docx()
    print("\nDone! Both files saved to:", REPORT_DIR)
